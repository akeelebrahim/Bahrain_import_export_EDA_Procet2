"""Builds the non-technical slide deck as PDF (required) and PPTX (editable).
Follows the course Presentation Standards: one chart per slide, chart fills the
slide, findings on the following slide, >=14pt, K/M/B number formatting, one
consistent high-contrast theme.
Run from the project root:  python scratch/build_presentation.py
Dsiclaimer: I used Ai and google search in parts on this code.
"""
import os
import textwrap
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import Rectangle, FancyBboxPatch
from matplotlib.backends.backend_pdf import PdfPages

ROOT = Path(r"~\Aqeel_Ebrahim_EDA_Project")
RAW = ROOT / "data" / "raw"
PRES = ROOT / "presentation"
ASSETS = PRES / "assets"
ASSETS.mkdir(parents=True, exist_ok=True)

# ----- Theme -----
NAVY = "#1F2A44"; GREY = "#4A4A4A"; BG = "#FFFFFF"
C_IMPORT = "#4C72B0"; C_EXPORT = "#DD8452"; C_POS = "#55A868"; C_NEG = "#C44E52"

# ============================================================ load + clean data
RENAME = {"export_value_bd": "value_bd", "export_weight_kg": "weight_kg",
          "import_value_bd": "value_bd", "import_weight_kg": "weight_kg",
          "qym_lwrdt_dynr_bhryny": "value_bd", "wzn_lwrdt_kjm": "weight_kg",
          "kmy_lwrdt": "quantity", "whd_lqys": "um",
          "export_value_usa": "value_usa", "export_quantity": "quantity",
          "import_value_usa": "value_usa", "import_quantity": "quantity",
          "qym_lwrdt_dwlr_mryky": "value_usa"}
DROP = ["n", "lshhr", "lsl", "ldwl"]

def load_detail(fname, flow, year_default=None):
    df = pd.read_csv(RAW / f"{fname}.csv", sep=";", dtype=str)
    df = df.rename(columns=RENAME).drop(columns=[c for c in DROP if c in df.columns])
    if "year" not in df.columns:
        df["year"] = year_default
    df["month_num"] = df["month"].str.extract(r"(\d+)").astype(int)
    for c in ["value_bd", "weight_kg"]:
        df[c] = pd.to_numeric(df[c], errors="coerce")
    df["commodity"] = df["commodity"].str.strip().str.title()
    df["country_name"] = df["country_name"].str.strip().str.title()
    df["commodity_no"] = df["commodity_no"].astype(str).str.strip()
    return df

imp23 = load_detail("imports_2023", "Import", 2023)
imp24 = load_detail("imports_2024", "Import", 2024)
texp24 = load_detail("total_exports_2024", "Total Export", 2024)
nexp24 = load_detail("national_exports_2024", "National Export", 2024)
rexp24 = load_detail("re_exports_2024", "Re-export", 2024)
annual = pd.read_csv(ROOT / "data" / "cleaned" / "annual_trade_clean.csv")

n_txn = len(imp24) + len(imp23) + len(texp24)
n_comm = pd.concat([imp24["commodity_no"], texp24["commodity_no"]]).nunique()
n_ctry = pd.concat([imp24["country_name"], texp24["country_name"]]).nunique()
imp_total = imp24["value_bd"].sum() / 1e6
exp_total = texp24["value_bd"].sum() / 1e6
nat_total = nexp24["value_bd"].sum() / 1e6
re_total = rexp24["value_bd"].sum() / 1e6
print(f"metrics: txn~{n_txn:,} comm={n_comm} ctry={n_ctry} imp={imp_total:.0f} exp={exp_total:.0f}")

# ============================================================ slide charts (titleless, 16:9)
plt.rcParams.update({"font.size": 16, "axes.labelsize": 16, "xtick.labelsize": 14,
                     "ytick.labelsize": 14, "legend.fontsize": 15, "axes.titlesize": 17,
                     "axes.titleweight": "bold", "figure.facecolor": BG, "axes.facecolor": BG,
                     "savefig.facecolor": BG})
CW, CH = 12.2, 5.35

def short(labels, n=34):
    return [l if len(l) <= n else l[:n-1] + "…" for l in labels]

# 1) annual trend
fig, ax = plt.subplots(figsize=(CW, CH))
ax.plot(annual["year"], annual["Total Imports"], marker="o", color=C_IMPORT, lw=2.5, label="Imports")
ax.plot(annual["year"], annual["Total Exports"], marker="o", color=C_EXPORT, lw=2.5, label="Exports")
bal = annual["Trade Balance"]
ax.bar(annual["year"], bal, color=[C_POS if v >= 0 else C_NEG for v in bal], alpha=0.35, label="Balance")
ax.axhline(0, color="grey", lw=0.8)
ax.set_ylabel("Million BD"); ax.legend(loc="upper left", ncol=3)
ax.yaxis.set_major_formatter(mticker.StrMethodFormatter("{x:,.0f}"))
fig.tight_layout(); fig.savefig(ASSETS / "annual.png", dpi=150); plt.close(fig)

# 2) balance + composition
fig, axes = plt.subplots(1, 2, figsize=(CW, CH))
axes[0].bar(["Imports", "Exports"], [imp_total, exp_total], color=[C_IMPORT, C_EXPORT])
for i, v in enumerate([imp_total, exp_total]):
    axes[0].text(i, v + 60, f"{v:,.0f}M", ha="center", fontweight="bold", fontsize=15)
axes[0].set_ylabel("Million BD"); axes[0].set_title("Non-oil imports vs exports (2024)")
axes[0].margins(y=0.15)
axes[1].bar(["National-origin", "Re-exports"], [nat_total, re_total], color=[C_EXPORT, "#E8B87F"])
for i, v in enumerate([nat_total, re_total]):
    axes[1].text(i, v + 40, f"{v/exp_total*100:.0f}%", ha="center", fontweight="bold", fontsize=15)
axes[1].set_ylabel("Million BD"); axes[1].set_title("Where exports come from"); axes[1].margins(y=0.15)
fig.tight_layout(); fig.savefig(ASSETS / "balance.png", dpi=150); plt.close(fig)

# 3) partners
ti = imp24.groupby("country_name")["value_bd"].sum().sort_values(ascending=False).head(10) / 1e6
te = texp24.groupby("country_name")["value_bd"].sum().sort_values(ascending=False).head(10) / 1e6
fig, axes = plt.subplots(1, 2, figsize=(CW, CH))
axes[0].barh(ti.index[::-1], ti.values[::-1], color=C_IMPORT); axes[0].set_title("Top import sources")
axes[1].barh(te.index[::-1], te.values[::-1], color=C_EXPORT); axes[1].set_title("Top export destinations")
for ax in axes:
    ax.set_xlabel("Million BD")
fig.tight_layout(); fig.savefig(ASSETS / "partners.png", dpi=150); plt.close(fig)

# 4) commodities
ci = imp24.groupby("commodity")["value_bd"].sum().sort_values(ascending=False).head(8) / 1e6
ce = texp24.groupby("commodity")["value_bd"].sum().sort_values(ascending=False).head(8) / 1e6
fig, axes = plt.subplots(1, 2, figsize=(CW, CH))
axes[0].barh(short(ci.index[::-1]), ci.values[::-1], color=C_IMPORT); axes[0].set_title("Top imports")
axes[1].barh(short(ce.index[::-1]), ce.values[::-1], color=C_EXPORT); axes[1].set_title("Top exports")
for ax in axes:
    ax.set_xlabel("Million BD")
fig.tight_layout(); fig.savefig(ASSETS / "commodities.png", dpi=150); plt.close(fig)

# 5) concentration
comm = imp24.groupby("commodity")["value_bd"].sum().sort_values(ascending=False)
cum = comm.cumsum() / comm.sum()
n80 = int((cum <= 0.80).sum() + 1)
fig, ax = plt.subplots(figsize=(CW, CH))
ax.plot(np.arange(1, len(cum) + 1), cum.values * 100, color=C_IMPORT, lw=3)
ax.axhline(80, color=C_NEG, ls="--"); ax.axvline(n80, color=C_NEG, ls="--")
ax.set_xlabel("Commodities ranked by value"); ax.set_ylabel("Cumulative % of import value")
ax.annotate(f"{n80} products = 80%", xy=(n80, 80), xytext=(n80 + 800, 55),
            fontsize=15, fontweight="bold", color=C_NEG,
            arrowprops=dict(arrowstyle="->", color=C_NEG))
fig.tight_layout(); fig.savefig(ASSETS / "concentration.png", dpi=150); plt.close(fig)
print("slide charts generated")

# ============================================================ slide content
SLIDES = [
    {"type": "cover", "title": "Buying Raw, Selling Refined",
     "subtitle": "Bahrain's Non-Oil Foreign Trade — an Exploratory Data Analysis",
     "author": "Aqeel Ebrahim  ·  Data Science DSB-1",
     "intro": "Exploring 760,000+ official trade records to see what Bahrain buys, what it makes and sells,\nand where its economic-diversification story really stands."},
    {"type": "text", "title": "The Problem", "bullets": [
        "Bahrain reports a headline trade surplus — but it is driven by oil.",
        "Strip oil out and 2024 flips to a ~1,190 million BD deficit.",
        "To support diversification under Vision 2030, we need to know:",
        "   –  What does Bahrain import and export (non-oil)?",
        "   –  Who are its main trading partners?",
        "   –  How concentrated, seasonal and changing is this trade?"]},
    {"type": "text", "title": "Our Approach", "bullets": [
        "Combined 6 official datasets from the Bahrain Open Data Portal.",
        "Cleaned inconsistent formats, duplicate columns and data-entry errors.",
        "Analysed value, partners and commodities — with charts for each question.",
        "Focus: the non-oil economy Bahrain is trying to grow."]},
    {"type": "metrics", "title": "The Data at a Glance",
     "metrics": [(f"{n_txn/1000:,.0f}K", "trade transactions"),
                 (f"{n_comm/1000:,.1f}K", "commodity lines (HS)"),
                 (f"{n_ctry}", "partner countries")],
     "footer": "Coverage:  2010–2023 (annual)  ·  2023–2024 (transaction detail)      Location:  Kingdom of Bahrain"},
    {"type": "text", "title": "Key Terms", "bullets": [
        "Non-oil trade  —  all trade except crude oil & refined petroleum.",
        "National-origin export  —  goods actually made in Bahrain.",
        "Re-export  —  imported goods sold on with little change.",
        "Trade balance  —  exports minus imports (surplus if positive).",
        "HS code  —  the international product-classification number."]},
    {"type": "chart", "title": "Trade swung from deficit to record surplus", "image": "annual.png"},
    {"type": "text", "title": "What the trend shows", "bullets": [
        "Five straight deficit years (2016–2020), worst in 2018 (–960M BD).",
        "Both flows collapsed in the 2020 pandemic, then rebounded hard.",
        "Record +2,516M BD surplus in 2022 — but largely an oil-price story.",
        "So we isolate the non-oil economy for the rest of the analysis."]},
    {"type": "chart", "title": "Without oil, Bahrain runs a trade deficit", "image": "balance.png"},
    {"type": "text", "title": "The hidden deficit", "bullets": [
        "2024 non-oil: imported ~5,872M BD, exported ~4,682M BD.",
        "That is a non-oil trade deficit of about 1,190 million BD.",
        "But ~83% of exports are national-origin — genuinely made in Bahrain.",
        "The export base reflects real production, not just re-selling."]},
    {"type": "chart", "title": "Exports lean heavily on Gulf neighbours", "image": "partners.png"},
    {"type": "text", "title": "Trading partners", "bullets": [
        "Saudi Arabia and the UAE dominate exports; top-5 = 54% of exports.",
        "Imports are more spread out (top-5 = 46%).",
        "Import leaders — China, Australia, Brazil — supply raw materials.",
        "Heavy reliance on a few export markets is a concentration risk."]},
    {"type": "chart", "title": "Bahrain buys raw materials and sells refined metal", "image": "commodities.png"},
    {"type": "text", "title": "A metals value-chain", "bullets": [
        "Biggest imports: iron ore and alumina (raw feedstock).",
        "Biggest exports: unwrought aluminium and iron/steel products.",
        "Bahrain refines imported raw materials using cheap energy.",
        "This value-added manufacturing is the core of diversification."]},
    {"type": "chart", "title": "A handful of products drive most trade", "image": "concentration.png"},
    {"type": "text", "title": "Recommendations", "bullets": [
        "Cut the non-oil deficit via import-substitution (iron ore, alumina, vehicles, electronics).",
        "Deepen the aluminium & steel value-chain — Bahrain's proven export engine.",
        "Diversify export markets beyond Saudi Arabia & the UAE (top-5 = 54%).",
        "Track the ~7% of products that drive 80% of trade value as an early warning."]},
    {"type": "text", "title": "Next Steps  ·  Thank You", "bullets": [
        "Extend to a full 2020–2026 multi-year panel.",
        "Add oil trade to reconcile with headline totals.",
        "Adjust for inflation to compare real trade volumes.",
        "Questions?  —  Aqeel Ebrahim, Data Science DSB-1"]},
]

# ============================================================ PDF renderer
def bg_ax(fig):
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off"); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.add_patch(Rectangle((0, 0), 1, 1, color=BG, zorder=-10))
    return ax

def title_block(ax, title):
    ax.add_patch(Rectangle((0.05, 0.88), 0.025, 0.06, color=C_EXPORT))  # accent
    ax.text(0.095, 0.905, title, fontsize=30, fontweight="bold", color=NAVY, va="center")

def bullets_block(ax, bullets, y0=0.76, size=19):
    """Render bullets with automatic line-wrapping so long text never overflows."""
    y = y0
    for b in bullets:
        stripped = b.strip()
        is_sub = stripped.startswith("–")
        x = 0.12 if is_sub else 0.09
        prefix = "" if is_sub else "•  "
        lines = textwrap.wrap(stripped, width=76) or [stripped]
        for j, ln in enumerate(lines):
            text = (prefix + ln) if j == 0 else ("   " + ln)
            ax.text(x, y, text, fontsize=size, color=GREY, va="top")
            y -= 0.058
        y -= 0.046

with PdfPages(PRES / "Aqeel_Ebrahim_EDA_Project_Presentation.pdf") as pdf:
    for s in SLIDES:
        fig = plt.figure(figsize=(13.333, 7.5)); ax = bg_ax(fig)
        if s["type"] == "cover":
            ax.add_patch(Rectangle((0, 0.60), 1, 0.02, color=C_EXPORT))
            ax.text(0.06, 0.74, s["title"], fontsize=52, fontweight="bold", color=NAVY)
            ax.text(0.06, 0.66, s["subtitle"], fontsize=22, color=GREY)
            ax.text(0.06, 0.40, s["intro"], fontsize=19, color=GREY)
            ax.text(0.06, 0.12, s["author"], fontsize=20, fontweight="bold", color=C_IMPORT)
        elif s["type"] == "text":
            title_block(ax, s["title"]); bullets_block(ax, s["bullets"])
        elif s["type"] == "metrics":
            title_block(ax, s["title"])
            xs = [0.20, 0.50, 0.80]
            for x, (big, lab) in zip(xs, s["metrics"]):
                ax.text(x, 0.55, big, fontsize=64, fontweight="bold", color=C_IMPORT, ha="center")
                ax.text(x, 0.42, lab, fontsize=20, color=GREY, ha="center")
            ax.text(0.5, 0.18, s["footer"], fontsize=17, color=NAVY, ha="center")
        elif s["type"] == "chart":
            title_block(ax, s["title"])
            img = plt.imread(ASSETS / s["image"])
            iax = fig.add_axes([0.045, 0.04, 0.91, 0.78]); iax.axis("off")
            iax.imshow(img, aspect="auto")
        pdf.savefig(fig); plt.close(fig)
print("PDF written")

# ============================================================ PPTX renderer
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(h): return RGBColor.from_string(h.lstrip("#"))
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY_R, GREY_R, IMP_R, EXP_R = rgb(NAVY), rgb(GREY), rgb(C_IMPORT), rgb(C_EXPORT)

    def add_title(slide, text):
        bar = slide.shapes.add_shape(1, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = EXP_R; bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.45), Inches(11.8), Inches(0.9)).text_frame
        tb.text = text; p = tb.paragraphs[0].runs[0].font
        p.size = Pt(30); p.bold = True; p.color.rgb = NAVY_R

    def add_bullets(slide, bullets, top=1.7):
        tf = slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.6), Inches(5.2)).text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("•  " + b) if not b.strip().startswith("–") else b
            p.space_after = Pt(12)
            p.runs[0].font.size = Pt(19); p.runs[0].font.color.rgb = GREY_R

    for s in SLIDES:
        sl = prs.slides.add_slide(blank)
        bgfill = sl.background.fill; bgfill.solid(); bgfill.fore_color.rgb = rgb(BG)
        if s["type"] == "cover":
            t = sl.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.4)).text_frame
            t.word_wrap = True; t.text = s["title"]
            r = t.paragraphs[0].runs[0].font; r.size = Pt(50); r.bold = True; r.color.rgb = NAVY_R
            for txt, sz, col, top in [(s["subtitle"], 22, GREY_R, 3.5), (s["intro"], 18, GREY_R, 4.6),
                                      (s["author"], 20, IMP_R, 6.4)]:
                b = sl.shapes.add_textbox(Inches(0.85), Inches(top), Inches(11.6), Inches(1.1)).text_frame
                b.word_wrap = True; b.text = txt
                b.paragraphs[0].runs[0].font.size = Pt(sz); b.paragraphs[0].runs[0].font.color.rgb = col
                if top == 6.4:
                    b.paragraphs[0].runs[0].font.bold = True
        elif s["type"] == "text":
            add_title(sl, s["title"]); add_bullets(sl, s["bullets"])
        elif s["type"] == "metrics":
            add_title(sl, s["title"])
            for j, (big, lab) in enumerate(s["metrics"]):
                x = Inches(0.9 + j * 4.0)
                b = sl.shapes.add_textbox(x, Inches(2.4), Inches(3.8), Inches(1.4)).text_frame
                b.text = big; b.paragraphs[0].alignment = PP_ALIGN.CENTER
                b.paragraphs[0].runs[0].font.size = Pt(60); b.paragraphs[0].runs[0].font.bold = True
                b.paragraphs[0].runs[0].font.color.rgb = IMP_R
                l = sl.shapes.add_textbox(x, Inches(3.9), Inches(3.8), Inches(0.8)).text_frame
                l.text = lab; l.paragraphs[0].alignment = PP_ALIGN.CENTER
                l.paragraphs[0].runs[0].font.size = Pt(19); l.paragraphs[0].runs[0].font.color.rgb = GREY_R
            f = sl.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.9)).text_frame
            f.word_wrap = True; f.text = s["footer"]
            f.paragraphs[0].alignment = PP_ALIGN.CENTER
            f.paragraphs[0].runs[0].font.size = Pt(16); f.paragraphs[0].runs[0].font.color.rgb = NAVY_R
        elif s["type"] == "chart":
            add_title(sl, s["title"])
            sl.shapes.add_picture(str(ASSETS / s["image"]), Inches(0.55), Inches(1.45), height=Inches(5.7))
    prs.save(str(PRES / "Aqeel_Ebrahim_EDA_Project_Presentation.pptx"))
    print("PPTX written")
except Exception as e:
    print("PPTX skipped:", type(e).__name__, e)
